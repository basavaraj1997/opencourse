import io
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import PyPDF2
import pytesseract
from PIL import Image
from pdf2image import convert_from_path
from PIL import Image


def extract_text_from_native_pdf(pdf_path):
    text = ""
    with open(pdf_path, "rb") as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        print('text extraction from native pdf')
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num-1]
            text += page.extract_text()
            image = page.images
            i=0
            for i, img in enumerate(image):
                with io.BytesIO(img.data) as img_buffer:                                          
                     text += pytesseract.image_to_string(Image.open(img_buffer))
            print(text)
    return text
  
def extract_text_from_pdf(pdf_path):
    if pdf_path.lower().endswith(".pdf"):
        try:
            text = extract_text_from_native_pdf(pdf_path)
        except Exception:            
                text =''
    elif pdf_path.lower().endswith(".docx"):
        text = extract_text_from_word(pdf_path)
        #text = "The file is not a PDF."
        return text
    else:
        text=extract_text_from_powerpoint(pdf_path)
        return text
def extract_text_from_word(file_path):
    doc = Document(file_path)
    text = ""
    
    for paragraph in doc.paragraphs:
        text += paragraph.text + '\n'
    return text


def extract_text_from_powerpoint(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

pdf_path = "C:/DOCPDF/22.pptx"
extracted_text = extract_text_from_pdf(pdf_path)
print(extracted_text)